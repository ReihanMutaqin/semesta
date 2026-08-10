import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_telkom_presentation():
    prs = Presentation()
    # Set slide dimensions to 16:9 widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Telkom Corporate Palette
    TELKOM_RED = RGBColor(224, 0, 0)         # #E00000
    DARK_NAVY = RGBColor(15, 23, 42)          # #0F172A
    CARD_BG = RGBColor(248, 250, 252)         # #F8FAFC
    CARD_BORDER = RGBColor(226, 232, 240)     # #E2E8F0
    TEXT_DARK = RGBColor(30, 41, 59)          # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139)      # #64748B
    EMERALD_GREEN = RGBColor(16, 185, 129)    # #10B981
    CYAN_BLUE = RGBColor(6, 182, 212)         # #06B6D4
    PURPLE_ACC = RGBColor(168, 85, 247)       # #A855F7
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="LAPORAN ANALISIS DATA SEMESTA (01.01.2026 - 10.08.2026)"):
        # Header banner shape
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        banner.fill.solid()
        banner.fill.fore_color.rgb = DARK_NAVY
        banner.line.color.rgb = DARK_NAVY

        # Red accent bar on left of banner
        red_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(1.1))
        red_bar.fill.solid()
        red_bar.fill.fore_color.rgb = TELKOM_RED
        red_bar.line.color.rgb = TELKOM_RED

        # Category / Subheader
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TELKOM_RED

        # Title Text
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.38), Inches(10.5), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

        # Telkom Logo Badge on top right
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.2), Inches(0.25), Inches(1.8), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TELKOM_RED
        badge.line.color.rgb = TELKOM_RED
        tf_badge = badge.text_frame
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = "TELKOM INDONESIA"
        p_badge.font.size = Pt(11)
        p_badge.font.bold = True
        p_badge.font.color.rgb = WHITE
        p_badge.alignment = PP_ALIGN.CENTER

    def add_footer(slide):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Telkom Operations & Support System • Rentang Acuan Data: 01.01.2026 - 10.08.2026 • Mas Faqih Report"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 1: COVER SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_NAVY
    bg1.line.color.rgb = DARK_NAVY

    # Decorative Telkom Red Accent Shape
    dec1 = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(0), Inches(3.833), Inches(7.5))
    dec1.fill.solid()
    dec1.fill.fore_color.rgb = TELKOM_RED
    dec1.line.color.rgb = TELKOM_RED
    dec1.rotation = 180

    # Title & Subtitle Box
    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.5), Inches(3.8))
    tf = tbox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "TELKOM OPERATIONS REPORT"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = TELKOM_RED

    p1 = tf.add_paragraph()
    p1.text = "Analisis Data Semesta"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "Modoroso • PDA HSI • IndiHome"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(4)

    p3 = tf.add_paragraph()
    p3.text = "Acuan Rentang Data: 01.01.2026 s/d 10.08.2026 (Total 88.011 Order)"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = EMERALD_GREEN
    p3.space_before = Pt(12)

    p4 = tf.add_paragraph()
    p4.text = "Evaluasi Jumlah Order Semesta, Order Terlama, Rata-rata Durasi PS, & Performa 1 Bulan Kebelakang"
    p4.font.size = Pt(11)
    p4.font.color.rgb = TEXT_MUTED
    p4.space_before = Pt(8)

    # Author Metadata Box
    abox = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(8.5), Inches(1.2))
    atf = abox.text_frame
    ap1 = atf.paragraphs[0]
    ap1.text = "Disusun untuk: Mas Faqih & Management Operations"
    ap1.font.size = Pt(11)
    ap1.font.bold = True
    ap1.font.color.rgb = WHITE

    ap2 = atf.add_paragraph()
    ap2.text = "Tanggal Laporan: 10 Agustus 2026 | Periode Data Terverifikasi: 01.01.2026 - 10.08.2026"
    ap2.font.size = Pt(10)
    ap2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & 4 CORE QUESTIONS
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Ringkasan Eksekutif: Jawaban 4 Pertanyaan Utama Operations")
    add_footer(slide2)

    # 4 Large Metric Cards
    metrics = [
        ("TOTAL ORDER SEMESTA (01.01.2026 - 10.08.2026)", "88.011", "Order Semesta Terdaftar", "80.795 Order PS Complete (91,8%)", TELKOM_RED),
        ("PS 1 BULAN KEBELAKANG", "5.910", "Order Selesai (30 Hari Terakhir)", "Didominasi Tipe CREATE (2.725) & MODIFY (2.109)", EMERALD_GREEN),
        ("RATA-RATA DURASI PS", "0,94 Hari", "~22,5 Jam (Tipe CREATE)", "Rata-rata DISCONNECT: 0,44 Hari (10,6 Jam)", DARK_NAVY),
        ("ORDER PENDING TERLAMA", "220,9 Hari", "Order Pending Terlama (CREATE)", "Pending MODIFY: 217,9 Hari", TELKOM_RED)
    ]

    for i, (m_title, m_val, m_sub1, m_sub2, accent_col) in enumerate(metrics):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.5 + row * 2.6)

        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        acc = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(0.12))
        acc.fill.solid()
        acc.fill.fore_color.rgb = accent_col
        acc.line.color.rgb = accent_col

        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_top = Inches(0.25)

        p = tf_card.paragraphs[0]
        p.text = m_title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED

        p_val = tf_card.add_paragraph()
        p_val.text = m_val
        p_val.font.size = Pt(32)
        p_val.font.bold = True
        p_val.font.color.rgb = accent_col
        p_val.space_before = Pt(3)

        p_s1 = tf_card.add_paragraph()
        p_s1.text = m_sub1
        p_s1.font.size = Pt(11)
        p_s1.font.bold = True
        p_s1.font.color.rgb = TEXT_DARK
        p_s1.space_before = Pt(3)

        p_s2 = tf_card.add_paragraph()
        p_s2.text = m_sub2
        p_s2.font.size = Pt(10)
        p_s2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: RINCIAN UTAMA PER TIPE TRANSAKSI (TABLE)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Rincian Metrik Utama per Tipe Transaksi (CRM Order Type)")
    add_footer(slide3)

    rows, cols = 7, 7
    table_shape = slide3.shapes.add_table(rows, cols, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.2))
    table = table_shape.table

    col_widths = [Inches(2.2), Inches(1.5), Inches(1.5), Inches(2.2), Inches(2.2), Inches(1.6), Inches(1.0)]
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = width

    headers = ["Tipe Transaksi", "Total Order", "Total PS", "Rata-rata Durasi PS", "Order Pending Terlama", "PS Terlama", "PS 1 Bulan"]
    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    table_data = [
        ["CREATE / NEW INSTALL", "51.962", "48.233 (92,8%)", "0,94 Hari (~22,5 Jam)", "220,9 Hari", "143,3 Hari", "2.725"],
        ["MODIFY", "19.157", "18.788 (98,1%)", "1,33 Hari (~31,9 Jam)", "217,9 Hari", "116,1 Hari", "2.109"],
        ["DISCONNECT", "11.337", "11.100 (97,9%)", "0,44 Hari (~10,6 Jam)", "167,0 Hari", "75,2 Hari", "905"],
        ["SUSPEND", "1.516", "1.050 (69,3%)", "< 1 Jam (Instant)", "-", "0,0 Hari", "0"],
        ["MIGRATE", "1.193", "1.066 (89,4%)", "2,07 Hari (~49,8 Jam)", "217,4 Hari", "95,6 Hari", "167"],
        ["UNSPECIFIED / LAINNYA", "2.321", "37 (1,6%)", "3,29 Hari (~79,0 Jam)", "220,8 Hari", "7,4 Hari", "4"]
    ]

    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = CARD_BG
            else:
                cell.fill.fore_color.rgb = WHITE
            
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK
            if col_idx == 0:
                p.font.bold = True
            if col_idx > 0:
                p.alignment = PP_ALIGN.RIGHT if col_idx in [1, 2, 6] else PP_ALIGN.CENTER
            if col_idx == 4 and cell_value != "-":
                p.font.color.rgb = TELKOM_RED
                p.font.bold = True

    # ==========================================
    # SLIDE 4: DEDICATED SLIDE - JUMLAH PS 1 BULAN KEBELAKANG PER TIPE TRANSAKSI
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Rincian Jumlah PS 1 Bulan Kebelakang per Tipe Transaksi", "ANALISIS PERFORMA 30 HARI TERAKHIR")
    add_footer(slide4)

    # 4 Focused Cards for PS 1 Month Back
    ps_cards = [
        ("CREATE / PASANG BARU", "2.725 PS", "46,1% dari Total PS Bulanan", "Rata-rata Durasi PS: 0,94 Hari (22,5 Jam)", TELKOM_RED),
        ("MODIFY / UBAH PAKET", "2.109 PS", "35,7% dari Total PS Bulanan", "Rata-rata Durasi PS: 1,33 Hari (31,9 Jam)", DARK_NAVY),
        ("DISCONNECT / CABUT", "905 PS", "15,3% dari Total PS Bulanan", "Rata-rata Durasi PS: 0,44 Hari (10,6 Jam)", CYAN_BLUE),
        ("MIGRATE / MIGRASI", "167 PS", "2,8% dari Total PS Bulanan", "Rata-rata Durasi PS: 2,07 Hari (49,8 Jam)", PURPLE_ACC)
    ]

    for idx, (p_title, p_val, p_pct, p_avg, p_col) in enumerate(ps_cards):
        left = Inches(0.6 + idx * 3.1)
        top = Inches(1.5)

        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.9), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        hbar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.9), Inches(0.6))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = p_col
        hbar.line.color.rgb = p_col

        tf_h = hbar.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = p_title
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0.7)
        tf_c.margin_left = Inches(0.15)

        p_lbl = tf_c.paragraphs[0]
        p_lbl.text = "JUMLAH PS (30 HARI)"
        p_lbl.font.size = Pt(9)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = TEXT_MUTED

        p_v = tf_c.add_paragraph()
        p_v.text = p_val
        p_v.font.size = Pt(30)
        p_v.font.bold = True
        p_v.font.color.rgb = p_col
        p_v.space_before = Pt(3)

        p_pct_lbl = tf_c.add_paragraph()
        p_pct_lbl.text = p_pct
        p_pct_lbl.font.size = Pt(10)
        p_pct_lbl.font.bold = True
        p_pct_lbl.font.color.rgb = TEXT_DARK
        p_pct_lbl.space_before = Pt(6)

        p_a = tf_c.add_paragraph()
        p_a.text = p_avg
        p_a.font.size = Pt(9.5)
        p_a.font.color.rgb = TEXT_MUTED
        p_a.space_before = Pt(6)

    # Bottom Highlight Box summarizing total 5,910 PS
    summary_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.5), Inches(12.133), Inches(1.3))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = DARK_NAVY
    summary_box.line.color.rgb = DARK_NAVY

    tf_s = summary_box.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = Inches(0.3)
    tf_s.margin_top = Inches(0.2)

    p_s_head = tf_s.paragraphs[0]
    p_s_head.text = "INSIGHT UTAMA PERFORMA PS 1 BULAN KEBELAKANG (TOTAL: 5.910 PS)"
    p_s_head.font.size = Pt(12)
    p_s_head.font.bold = True
    p_s_head.font.color.rgb = TELKOM_RED

    p_s_desc = tf_s.add_paragraph()
    p_s_desc.text = "• Sebanyak 81,8% dari total PS bulanan disumbangkan oleh aktivitas Pasang Baru (CREATE: 2.725 PS) & Perubahan Paket (MODIFY: 2.109 PS).\n• Kecepatan penyelesaian PS sangat responsif dengan rata-rata SLA CREATE 0,94 Hari (~22,5 Jam) & DISCONNECT 0,44 Hari (~10,6 Jam)."
    p_s_desc.font.size = Pt(10.5)
    p_s_desc.font.color.rgb = WHITE
    p_s_desc.space_before = Pt(4)

    # ==========================================
    # SLIDE 5: BREAKDOWN SEGMENTASI (MODOROSO, PDA HSI, INDIHOME)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Breakdown Segmentasi Order: Modoroso, PDA HSI, & IndiHome")
    add_footer(slide5)

    seg_data = [
        ("MODOROSO", "26.058", "22.378 PS (85,9%)", "Rata-rata PS: 1,25 Hari", "Order Pending Terlama: 220,9 Hari", "PS 1 Bulan: 2.104", TELKOM_RED),
        ("PDA HSI", "39.308", "37.123 PS (94,4%)", "Rata-rata PS: 0,63 Hari", "Order Pending Terlama: 217,9 Hari", "PS 1 Bulan: 3.325", DARK_NAVY),
        ("INDIHOME", "5.900", "5.726 PS (97,1%)", "Rata-rata PS: 0,57 Hari", "Order Pending Terlama: 167,0 Hari", "PS 1 Bulan: 457", EMERALD_GREEN),
        ("ENTERPRISE / LAINNYA", "16.745", "15.568 PS (93,0%)", "Rata-rata PS: 1,89 Hari", "Order Pending Terlama: 220,8 Hari", "PS 1 Bulan: 24", TEXT_MUTED)
    ]

    for idx, (s_name, s_tot, s_ps, s_avg, s_max, s_m30, col_acc) in enumerate(seg_data):
        left = Inches(0.6 + idx * 3.1)
        top = Inches(1.5)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.9), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        h_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.9), Inches(0.8))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = col_acc
        h_box.line.color.rgb = col_acc

        tf_h = h_box.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = s_name
        p_h.font.size = Pt(13)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0.9)
        tf_c.margin_left = Inches(0.15)
        tf_c.margin_right = Inches(0.15)

        p_tot_label = tf_c.paragraphs[0]
        p_tot_label.text = "TOTAL ORDER"
        p_tot_label.font.size = Pt(9)
        p_tot_label.font.bold = True
        p_tot_label.font.color.rgb = TEXT_MUTED

        p_tot_val = tf_c.add_paragraph()
        p_tot_val.text = s_tot
        p_tot_val.font.size = Pt(28)
        p_tot_val.font.bold = True
        p_tot_val.font.color.rgb = TEXT_DARK
        p_tot_val.space_before = Pt(2)

        items = [
            ("Status PS", s_ps, EMERALD_GREEN),
            ("Rata-rata PS", s_avg, TEXT_DARK),
            ("Order Terlama", s_max, TELKOM_RED),
            ("PS 1 Bulan", s_m30, DARK_NAVY)
        ]

        for label, val_text, txt_col in items:
            p_lbl = tf_c.add_paragraph()
            p_lbl.text = label.upper()
            p_lbl.font.size = Pt(9)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = TEXT_MUTED
            p_lbl.space_before = Pt(12)

            p_val = tf_c.add_paragraph()
            p_val.text = val_text
            p_val.font.size = Pt(11)
            p_val.font.bold = True
            p_val.font.color.rgb = txt_col

    # ==========================================
    # SLIDE 6: ANALISIS ORDER TERLAMA & BOTTLENECK SLA
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Analisis Kendala Order Terlama & Bottleneck SLA Operations")
    add_footer(slide6)

    box_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = CARD_BG
    box_l.line.color.rgb = CARD_BORDER

    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.3)

    p = tf_l.paragraphs[0]
    p.text = "TEMUAN UTAMA ORDER TERLAMA (> 200 HARI)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TELKOM_RED

    findings = [
        "Order Terlama Pending Mencapai 220,9 Hari pada tipe transaksi CREATE dan 217,9 Hari pada tipe MODIFY.",
        "Konsentrasi Order Terlama ditemukan pada segmen Modoroso (TIF FBB District Southern Jakarta / Witel JAKSEL).",
        "Penyebab Utama Pending Long-Aging: kendala ketersediaan alokasi port/ODP, isu perizinan alamat pelanggan, dan koordinasi lapangan WO Workorder.",
        "Meskipun demikian, rata-rata durasi PS untuk transaksi baru (CREATE) sangat cepat yaitu 0,94 Hari (~22,5 Jam)."
    ]

    for item in findings:
        p_item = tf_l.add_paragraph()
        p_item.text = "• " + item
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = TEXT_DARK
        p_item.space_before = Pt(10)

    box_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = CARD_BG
    box_r.line.color.rgb = CARD_BORDER

    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.3)

    p_r = tf_r.paragraphs[0]
    p_r.text = "REKOMENDASI PERBAIKAN OPERASIONAL"
    p_r.font.size = Pt(13)
    p_r.font.bold = True
    p_r.font.color.rgb = DARK_NAVY

    recs = [
        "Pembersihan Backlog (Clearing Long-Aging): Membentuk Task Force khusus penanganan order berusia > 30 hari untuk validasi fisik ODP di Witel Jaksel & Jaktim.",
        "Otomasi Filter & Dashboard Monitoring: Menggunakan Dashboard Data Semesta Vercel dengan filter urutan 'Order Terlama' untuk penanganan harian teknisi.",
        "Standardisasi SLA Tipe Transaksi: Mempertahankan SLA CREATE < 24 jam dan mempercepat proses administrasi tipe MIGRATE (rata-rata 2,07 hari).",
        "Integrasi Sistem Berkala: Mengunggah file laporan bulanan .xlsx ke web dashboard untuk menjaga visibilitas real-time manajemen."
    ]

    for rec in recs:
        p_rec = tf_r.add_paragraph()
        p_rec.text = "✔ " + rec
        p_rec.font.size = Pt(11)
        p_rec.font.color.rgb = TEXT_DARK
        p_rec.space_before = Pt(10)

    # Save presentation to both paths
    paths = [
        r"C:\PROJEK\Data Semesta\Laporan_Analisis_Data_Semesta_Telkom_v3.pptx",
        r"C:\PROJEK\Data Semesta\Laporan_Analisis_Data_Semesta_Telkom_v2.pptx",
        r"C:\PROJEK\Data Semesta\Laporan_Analisis_Data_Semesta_Telkom.pptx"
    ]
    for output_path in paths:
        try:
            prs.save(output_path)
            print(f"Presentation saved successfully at: {output_path}")
        except Exception as e:
            print(f"Could not write to {output_path} (file might be open): {e}")

if __name__ == '__main__':
    create_telkom_presentation()
